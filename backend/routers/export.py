from fastapi import APIRouter, Depends, HTTPException, status
from fastapi.responses import StreamingResponse
from sqlalchemy.orm import Session
from database import get_db
from models import Project, Section, User
from auth import get_current_user
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as PptRGBColor
import io

router = APIRouter(prefix="/export", tags=["export"])

@router.get("/projects/{project_id}/download")
async def export_document(
    project_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """
    Export project as DOCX or PPTX file
    """
    # Get project with sections
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    
    if not project:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Project not found"
        )
    
    # Get sections ordered by order
    sections = db.query(Section).filter(
        Section.project_id == project_id
    ).order_by(Section.order).all()
    
    if not sections or not any(s.content for s in sections):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="No content to export. Generate content first."
        )
    
    try:
        # Create filename
        safe_title = "".join(c if c.isalnum() or c in (' ', '-', '_') else '_' for c in project.title)
        safe_title = safe_title.replace(' ', '_')
        
        if project.document_type.value == "docx":
            # Create DOCX in memory
            file_stream = create_docx(project, sections)
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            filename = f"{safe_title}.docx"
        else:  # pptx
            # Create PPTX in memory
            file_stream = create_pptx(project, sections)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            filename = f"{safe_title}.pptx"
        
        # Return as streaming response
        return StreamingResponse(
            file_stream,
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename={filename}",
                "Content-Type": media_type
            }
        )
    
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error exporting document: {str(e)}"
        )

def create_docx(project: Project, sections: list) -> io.BytesIO:
    """
    Create DOCX file in memory from project content with:
    - Page 1: Title + Description + Table of Contents (NO BLANK PAGE)
    - Page 2+: Content sections
    
    Returns: BytesIO stream
    """
    # Create document
    doc = Document()
    
    # Set default font
    style = doc.styles['Normal']
    style.font.name = 'Calibri'
    style.font.size = Pt(11)
    
    # ===== PAGE 1: TITLE + DESCRIPTION + TOC =====
    
    # Add vertical spacing before title (smaller)
    doc.add_paragraph()
    
    # Add title - centered, large, colored
    title = doc.add_heading(project.title, 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title.runs[0]
    title_run.font.color.rgb = RGBColor(31, 78, 121)
    title_run.font.size = Pt(36)
    title_run.font.bold = True
    
    # Small spacing
    spacing_para = doc.add_paragraph()
    spacing_para.paragraph_format.space_after = Pt(6)
    
    # Add topic as subtitle - centered, italic
    topic_para = doc.add_paragraph(project.topic)
    topic_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    topic_para_run = topic_para.runs[0]
    topic_para_run.italic = True
    topic_para_run.font.size = Pt(14)
    topic_para_run.font.color.rgb = RGBColor(89, 89, 89)
    
    # Moderate spacing before TOC
    spacing_para2 = doc.add_paragraph()
    spacing_para2.paragraph_format.space_after = Pt(12)
    
    # ===== TABLE OF CONTENTS ON SAME PAGE =====
    # Filter sections with content
    content_sections = [s for s in sections if s.content]
    
    # Add TOC heading
    toc_heading = doc.add_heading('Table of Contents', 2)
    toc_heading_run = toc_heading.runs[0]
    toc_heading_run.font.color.rgb = RGBColor(31, 78, 121)
    toc_heading_run.font.size = Pt(14)
    
    # Minimal spacing after TOC heading
    spacing_para3 = doc.add_paragraph()
    spacing_para3.paragraph_format.space_after = Pt(6)
    
    # Add TOC entries with page numbers (approximate)
    current_page = 2  # Content starts on page 2
    
    for idx, section in enumerate(content_sections):
        # Create TOC entry with dot leaders (compact)
        toc_entry = doc.add_paragraph()
        toc_entry.paragraph_format.left_indent = Inches(0.3)
        toc_entry.paragraph_format.first_line_indent = Inches(-0.3)
        toc_entry.paragraph_format.space_after = Pt(4)  # Tight spacing
        
        # Section title
        title_run = toc_entry.add_run(f"{idx + 1}. {section.title}")
        title_run.font.size = Pt(10)
        title_run.font.color.rgb = RGBColor(31, 78, 121)
        
        # Dots
        dots_run = toc_entry.add_run(" " + "." * 40)
        dots_run.font.size = Pt(10)
        dots_run.font.color.rgb = RGBColor(180, 180, 180)
        
        # Page number
        page_run = toc_entry.add_run(f" {current_page}")
        page_run.font.size = Pt(10)
        page_run.font.bold = True
        
        # Estimate pages (roughly 500 words = 1 page)
        word_count = len(section.content.split())
        estimated_pages = max(1, word_count // 500)
        current_page += estimated_pages
    
    # ===== PAGE BREAK AFTER TOC (NO BLANK PAGE) =====
    doc.add_page_break()
    
    # ===== CONTENT SECTIONS (PAGE 2+) =====
    for idx, section in enumerate(content_sections):
        # Section heading with number
        heading_text = f"{idx + 1}. {section.title}"
        heading = doc.add_heading(heading_text, 1)
        heading_run = heading.runs[0]
        heading_run.font.color.rgb = RGBColor(31, 78, 121)
        heading_run.font.size = Pt(22)
        
        # Minimal spacing after heading
        spacing = doc.add_paragraph()
        spacing.paragraph_format.space_after = Pt(6)
        
        # Section content - handle multi-paragraph content
        paragraphs = section.content.split('\n\n')
        
        for para_idx, para_text in enumerate(paragraphs):
            para_text = para_text.strip()
            if para_text:
                # Handle any remaining single newlines within paragraphs
                para_text = para_text.replace('\n', ' ')
                
                content_para = doc.add_paragraph(para_text)
                content_para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                content_para.paragraph_format.line_spacing = 1.15
                content_para.paragraph_format.space_after = Pt(10)
                content_para.paragraph_format.first_line_indent = Inches(0.5)
        
        # Add spacing between sections (but not after last one)
        if idx < len(content_sections) - 1:
            spacing_between = doc.add_paragraph()
            spacing_between.paragraph_format.space_after = Pt(6)
            doc.add_page_break()
    
    # Save to BytesIO stream
    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    
    return file_stream

def create_pptx(project: Project, sections: list) -> io.BytesIO:
    """
    Create PPTX file in memory from project content
    Returns: BytesIO stream
    """
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)
    
    # Define color scheme
    TITLE_COLOR = PptRGBColor(31, 78, 121)
    TEXT_COLOR = PptRGBColor(50, 50, 50)
    
    # ===== TITLE SLIDE =====
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PptRGBColor(245, 245, 245)
    
    # Add title text box (centered vertically and horizontally)
    left = PptInches(0.5)
    top = PptInches(2.5)
    width = PptInches(9)
    height = PptInches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    title_p = title_frame.paragraphs[0]
    title_p.text = project.title
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.size = PptPt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = TITLE_COLOR
    
    # Add subtitle text box
    subtitle_top = PptInches(4.2)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, PptInches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = project.topic
    subtitle_p.alignment = PP_ALIGN.CENTER
    subtitle_p.font.size = PptPt(24)
    subtitle_p.font.italic = True
    subtitle_p.font.color.rgb = PptRGBColor(100, 100, 100)
    
    # ===== CONTENT SLIDES =====
    for section in sections:
        if section.content:
            # Use blank layout for custom formatting
            content_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Add background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = PptRGBColor(255, 255, 255)
            
            # Add title
            title_left = PptInches(0.5)
            title_top = PptInches(0.4)
            title_width = PptInches(9)
            title_height = PptInches(0.8)
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            
            title_p = title_frame.paragraphs[0]
            title_p.text = section.title
            title_p.font.size = PptPt(40)
            title_p.font.bold = True
            title_p.font.color.rgb = TITLE_COLOR
            
            # Add content
            content_left = PptInches(0.7)
            content_top = PptInches(1.4)
            content_width = PptInches(8.6)
            content_height = PptInches(5.8)
            
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            # Split content into bullet points (already cleaned by gemini_service)
            bullet_lines = section.content.split('\n')
            
            first_bullet = True
            for line in bullet_lines:
                line = line.strip()
                if not line:
                    continue
                
                if first_bullet:
                    p = text_frame.paragraphs[0]
                    first_bullet = False
                else:
                    p = text_frame.add_paragraph()
                
                p.text = line
                p.level = 0
                p.font.size = PptPt(18)
                p.font.color.rgb = TEXT_COLOR
                p.space_before = PptPt(8)
                p.space_after = PptPt(8)
                p.line_spacing = 1.3
    
    # ===== CLOSING SLIDE =====
    closing_layout = prs.slide_layouts[6]
    closing_slide = prs.slides.add_slide(closing_layout)
    
    # Add background
    background = closing_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PptRGBColor(31, 78, 121)
    
    # Add closing text
    closing_left = PptInches(0.5)
    closing_top = PptInches(3)
    closing_width = PptInches(9)
    closing_height = PptInches(1.5)
    
    closing_box = closing_slide.shapes.add_textbox(closing_left, closing_top, closing_width, closing_height)
    closing_frame = closing_box.text_frame
    closing_frame.word_wrap = True
    closing_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    closing_p = closing_frame.paragraphs[0]
    closing_p.text = "Thank You"
    closing_p.alignment = PP_ALIGN.CENTER
    closing_p.font.size = PptPt(48)
    closing_p.font.bold = True
    closing_p.font.color.rgb = PptRGBColor(255, 255, 255)
    
    # Save to BytesIO stream
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    return file_stream
